"""PowerPoint report generation using python-pptx."""

import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BRAND_BLUE = RGBColor(0x2E, 0x50, 0x90)
BRAND_ORANGE = RGBColor(0xE8, 0x83, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)


def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=12,
                  bold=False, color=DARK_GREY, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txbox


def _add_table(slide, df, left, top, width, height):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    col_width = width // cols
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for j, col_name in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_BLUE
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(9)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(8)
                para.font.color.rgb = DARK_GREY
                para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def _title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, BRAND_BLUE)
    _add_text_box(
        slide, Inches(1), Inches(2.5), Inches(8), Inches(1.2),
        "Quarterly Investment Report",
        font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, Inches(1), Inches(3.8), Inches(8), Inches(0.6),
        "Q4 2024 — Sample Fund",
        font_size=18, color=RGBColor(0xBB, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER,
    )


def _section_slide(prs, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_BLUE)
    _add_text_box(
        slide, Inches(1), Inches(3), Inches(8), Inches(1),
        title, font_size=28, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
    )


def _chart_slide(prs, title: str, chart_path: str | io.BytesIO, commentary: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
        title, font_size=22, bold=True, color=BRAND_BLUE,
    )
    if isinstance(chart_path, io.BytesIO):
        chart_path.seek(0)
    slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.1), Inches(8), Inches(4.5))
    if commentary:
        _add_text_box(
            slide, Inches(0.8), Inches(5.8), Inches(8.2), Inches(1.2),
            commentary, font_size=10, color=DARK_GREY,
        )


def _table_slide(prs, title: str, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
        title, font_size=22, bold=True, color=BRAND_BLUE,
    )
    _add_table(slide, df, Inches(0.4), Inches(1.2), Inches(9.2), Inches(4.5))


def _text_slide(prs, title: str, body: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
        title, font_size=22, bold=True, color=BRAND_BLUE,
    )
    _add_text_box(
        slide, Inches(0.8), Inches(1.2), Inches(8.2), Inches(5.5),
        body, font_size=13, color=DARK_GREY,
    )


def generate_pptx(
    charts: dict[str, str | io.BytesIO],
    holdings_df,
    commentary: dict,
    output_path: str = "output/report.pptx",
) -> str:
    """Build the full PowerPoint report."""
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    _title_slide(prs)

    # Slide 2: Executive Summary
    _text_slide(prs, "Executive Summary", commentary["executive_summary"])

    # Slide 3: Section divider
    _section_slide(prs, "Performance")

    # Slide 4: Line Chart
    _chart_slide(
        prs, "Portfolio vs Benchmark Performance", charts["line_chart"],
        commentary.get(
            "performance_commentary",
            "The portfolio outperformed the benchmark by 180bps over the trailing period.",
        )[:300],
    )

    # Slide 5: Histogram
    _chart_slide(
        prs, "Return Distribution", charts["histogram"],
        commentary.get(
            "return_distribution_commentary",
            "Monthly returns exhibit positive skew consistent with quality tilt.",
        )[:300],
    )

    # Slide 6: Section divider
    _section_slide(prs, "Asset Allocation & Holdings")

    # Slide 7: Stacked Area
    _chart_slide(
        prs, "Asset Allocation Over Time", charts["stacked_area"],
        commentary.get(
            "allocation_commentary",
            "Allocation shifted toward equities in response to improving economic outlook.",
        )[:300],
    )

    # Slide 8: Holdings Table
    _table_slide(prs, "Top Holdings", holdings_df)

    # Slide 9: Section divider
    _section_slide(prs, "Risk & Outlook")

    # Slide 10: Spider Chart
    _chart_slide(
        prs, "Risk Factor Comparison", charts["spider_chart"],
        commentary.get(
            "risk_assessment",
            "Fund A shows higher market and liquidity risk; Fund B tilts toward credit exposure.",
        )[:300],
    )

    # Slide 11: Risk Assessment
    _text_slide(prs, "Risk Assessment", commentary["risk_assessment"])

    # Slide 12: Market Outlook
    _text_slide(prs, "Market Outlook", commentary["market_outlook"])

    prs.save(output_path)
    return output_path
