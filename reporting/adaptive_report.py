"""Adaptive report builder — executes a plan produced by the LLM planner.

All builder functions accept an optional `theme` dict. When provided, colors,
fonts, logos, and branding text are taken from the theme. When None, the
built-in defaults are used (backward-compatible).
"""

import io
from pathlib import Path

from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm, mm
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from chart_library import CHART_REGISTRY
from theme import DEFAULT_THEME

WIDTH, HEIGHT = A4


# ─────────────────────────────────────────────────────────────────────
# Helpers — extract ReportLab / PPTX color objects from theme
# ─────────────────────────────────────────────────────────────────────

def _rl(hex_color: str):
    """Convert hex string to ReportLab color."""
    return colors.HexColor(hex_color)


def _rgb(hex_color: str) -> RGBColor:
    """Convert hex string to python-pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _c(theme: dict) -> dict:
    """Shorthand to get the theme color dict."""
    return theme["colors"]


# ─────────────────────────────────────────────────────────────────────
# Chart generation — execute only the charts the planner selected
# ─────────────────────────────────────────────────────────────────────

def generate_planned_charts(plan: dict, data: dict, output_dir: str) -> dict[str, str]:
    """Generate only the charts specified in the plan."""
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    chart_paths = {}
    for section in plan["sections"]:
        chart_id = section.get("chart")
        if chart_id and chart_id in CHART_REGISTRY:
            path = f"{output_dir}/{chart_id}.png"
            fn = CHART_REGISTRY[chart_id]["fn"]
            chart_paths[chart_id] = fn(data, path)
    return chart_paths


# ─────────────────────────────────────────────────────────────────────
# PDF Builder
# ─────────────────────────────────────────────────────────────────────

def _pdf_styles(theme: dict) -> dict:
    ss = getSampleStyleSheet()
    primary = _rl(_c(theme)["primary"])
    heading_font = theme["fonts"]["heading"]
    body_font = theme["fonts"]["body"]

    return {
        "title": ParagraphStyle(
            "RT", parent=ss["Title"], fontSize=28, leading=34,
            textColor=primary, spaceAfter=6*mm, alignment=TA_CENTER,
            fontName=heading_font,
        ),
        "subtitle": ParagraphStyle(
            "RS", parent=ss["Normal"], fontSize=12, leading=16,
            textColor=colors.grey, spaceAfter=8*mm, alignment=TA_CENTER,
            fontName=body_font,
        ),
        "client_name": ParagraphStyle(
            "CN", parent=ss["Normal"], fontSize=10, leading=14,
            textColor=_rl(_c(theme)["secondary"]), spaceAfter=15*mm,
            alignment=TA_CENTER, fontName=body_font,
        ),
        "character": ParagraphStyle(
            "RC", parent=ss["Normal"], fontSize=11, leading=15,
            textColor=colors.HexColor("#666666"), spaceAfter=10*mm,
            alignment=TA_CENTER, fontName=f"{body_font}-Oblique"
            if body_font == "Helvetica" else body_font,
        ),
        "heading": ParagraphStyle(
            "RH", parent=ss["Heading1"], fontSize=16, leading=20,
            textColor=primary, spaceBefore=10*mm, spaceAfter=4*mm,
            fontName=heading_font,
        ),
        "body": ParagraphStyle(
            "RB", parent=ss["BodyText"], fontSize=10, leading=14,
            alignment=TA_JUSTIFY, spaceAfter=3*mm, fontName=body_font,
        ),
        "table_header": ParagraphStyle(
            "TH", parent=ss["Normal"], fontSize=9, leading=12,
            textColor=colors.white, alignment=TA_CENTER, fontName=heading_font,
        ),
        "table_cell": ParagraphStyle(
            "TC", parent=ss["Normal"], fontSize=9, leading=12,
            alignment=TA_CENTER, fontName=body_font,
        ),
        "table_cell_left": ParagraphStyle(
            "TCL", parent=ss["Normal"], fontSize=9, leading=12,
            alignment=TA_LEFT, fontName=body_font,
        ),
    }


def _make_header_footer(theme: dict):
    """Return a header/footer callback bound to the theme."""
    primary = _rl(_c(theme)["primary"])
    header_text = theme["pdf"]["header_text"]
    footer_text = theme["pdf"]["footer_text"]
    logo_path = theme.get("logo")
    body_font = theme["fonts"]["body"]

    def _header_footer(canvas, doc):
        canvas.saveState()
        # Header line
        canvas.setStrokeColor(primary)
        canvas.setLineWidth(1.5)
        canvas.line(2*cm, HEIGHT - 1.5*cm, WIDTH - 2*cm, HEIGHT - 1.5*cm)
        canvas.setFont(body_font, 8)
        canvas.setFillColor(colors.grey)
        canvas.drawString(2*cm, HEIGHT - 1.3*cm, header_text)

        # Logo in header (right side)
        if logo_path and Path(logo_path).exists():
            canvas.drawImage(
                logo_path,
                WIDTH - 4.5*cm, HEIGHT - 1.45*cm, width=2.2*cm, height=0.7*cm,
                preserveAspectRatio=True, mask="auto",
            )

        # Footer
        canvas.line(2*cm, 1.8*cm, WIDTH - 2*cm, 1.8*cm)
        canvas.drawString(2*cm, 1.2*cm, footer_text)
        canvas.drawRightString(WIDTH - 2*cm, 1.2*cm, f"Page {doc.page}")
        canvas.restoreState()

    return _header_footer


def _build_holdings_table(df, styles: dict, theme: dict) -> Table:
    primary = _rl(_c(theme)["primary"])
    bg_alt = _rl(_c(theme)["background_alt"])
    grid_color = _rl(_c(theme)["grid"])

    header = [Paragraph(str(c), styles["table_header"]) for c in df.columns]
    rows = [header]
    for _, row in df.iterrows():
        cells = []
        for i, val in enumerate(row):
            s = styles["table_cell_left"] if i == 0 else styles["table_cell"]
            cells.append(Paragraph(str(val), s))
        rows.append(cells)
    n_cols = len(df.columns)
    col_w = (WIDTH - 4*cm) / n_cols
    tbl = Table(rows, colWidths=[col_w]*n_cols, repeatRows=1)
    tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), primary),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("ALIGN", (0, 1), (0, -1), "LEFT"),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, bg_alt]),
        ("GRID", (0, 0), (-1, -1), 0.5, grid_color),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]))
    return tbl


def build_pdf(plan: dict, chart_paths: dict, data: dict, output_path: str,
              theme: dict | None = None) -> str:
    """Build a PDF from a plan, themed by the given theme dict."""
    theme = theme or DEFAULT_THEME
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    styles = _pdf_styles(theme)

    frame = Frame(2*cm, 2.5*cm, WIDTH - 4*cm, HEIGHT - 4.5*cm, id="main")
    hf = _make_header_footer(theme)
    template = PageTemplate(id="std", frames=frame, onPage=hf)
    doc = BaseDocTemplate(output_path, pagesize=A4, pageTemplates=[template])

    story = []

    # ── Cover ────────────────────────────────────────────────────
    # Logo on cover
    logo_path = theme.get("logo")
    if logo_path and Path(logo_path).exists():
        story.append(Spacer(1, 2*cm))
        logo = Image(logo_path, width=5*cm, height=1.7*cm)
        logo.hAlign = "CENTER"
        story.append(logo)
        story.append(Spacer(1, 2*cm))
    else:
        story.append(Spacer(1, 5*cm))

    story.append(Paragraph(plan.get("report_title", "Investment Report"), styles["title"]))
    story.append(Paragraph(plan.get("report_subtitle", ""), styles["subtitle"]))

    if theme.get("client_name"):
        story.append(Paragraph(theme["client_name"], styles["client_name"]))

    if plan.get("fund_character"):
        story.append(Paragraph(f'"{plan["fund_character"]}"', styles["character"]))
    story.append(PageBreak())

    # ── Sections ─────────────────────────────────────────────────
    for i, section in enumerate(plan["sections"]):
        story.append(Paragraph(section["title"], styles["heading"]))
        story.append(Paragraph(section["commentary"], styles["body"]))

        chart_id = section.get("chart")
        if chart_id and chart_id in chart_paths:
            img_w = 12*cm if chart_id in ("spider_chart", "pie_chart") else 14*cm
            img = Image(chart_paths[chart_id], width=img_w, height=img_w * 0.57)
            story.append(img)
            story.append(Spacer(1, 3*mm))

        if section.get("include_holdings_table") or (
            chart_id in ("contributor_bar", "pie_chart") and
            "holding" in section.get("title", "").lower()
        ):
            story.append(Spacer(1, 3*mm))
            story.append(_build_holdings_table(data["holdings"], styles, theme))

        if chart_id and i < len(plan["sections"]) - 1:
            story.append(PageBreak())

    doc.build(story)
    return output_path


# ─────────────────────────────────────────────────────────────────────
# PPTX Builder
# ─────────────────────────────────────────────────────────────────────

def _pptx_text(slide, left, top, width, height, text, size=12,
               bold=False, color=None, align=PP_ALIGN.LEFT, font_name=None):
    color = color or RGBColor(0x33, 0x33, 0x33)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if font_name:
        p.font.name = font_name


def _pptx_table(slide, df, left, top, width, height, theme: dict):
    primary = _rgb(_c(theme)["primary"])
    bg_alt = _rgb(_c(theme)["background_alt"])
    text_dark = _rgb(_c(theme)["text_dark"])

    rows, cols = df.shape[0] + 1, df.shape[1]
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    cw = width // cols
    for i in range(cols):
        tbl.columns[i].width = cw
    for j, c in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.text = str(c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (_, row) in enumerate(df.iterrows()):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_alt
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(7)
                p.font.color.rgb = text_dark
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def build_pptx(plan: dict, chart_paths: dict, data: dict, output_path: str,
               theme: dict | None = None) -> str:
    """Build a PPTX from a plan, themed by the given theme dict."""
    theme = theme or DEFAULT_THEME
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    primary = _rgb(_c(theme)["primary"])
    secondary = _rgb(_c(theme)["secondary"])
    text_dark = _rgb(_c(theme)["text_dark"])
    title_bg = _rgb(theme["pptx"]["title_slide_bg"])
    section_bg = _rgb(theme["pptx"]["section_slide_bg"])
    white = RGBColor(0xFF, 0xFF, 0xFF)
    subtitle_color = RGBColor(
        min(255, title_bg.red + 100),
        min(255, title_bg.green + 100),
        min(255, title_bg.blue + 60),
    )

    # ── Title slide ──────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = title_bg

    # Logo on title slide
    logo_path = theme.get("logo")
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(
            logo_path, Inches(3.5), Inches(1.0), Inches(3), Inches(1),
        )
        title_top = Inches(2.3)
    else:
        title_top = Inches(2.2)

    _pptx_text(slide, Inches(1), title_top, Inches(8), Inches(1.2),
               plan.get("report_title", "Report"), size=36, bold=True,
               color=white, align=PP_ALIGN.CENTER)
    _pptx_text(slide, Inches(1), Inches(title_top.inches + 1.3), Inches(8), Inches(0.6),
               plan.get("report_subtitle", ""), size=18,
               color=subtitle_color, align=PP_ALIGN.CENTER)
    if theme.get("client_name"):
        _pptx_text(slide, Inches(1), Inches(title_top.inches + 2.1), Inches(8), Inches(0.5),
                   theme["client_name"], size=14,
                   color=subtitle_color, align=PP_ALIGN.CENTER)

    # ── Sections ─────────────────────────────────────────────────
    for section in plan["sections"]:
        chart_id = section.get("chart")

        if chart_id and chart_id in chart_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_text(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                       section["title"], size=22, bold=True, color=primary)
            slide.shapes.add_picture(
                chart_paths[chart_id],
                Inches(0.8), Inches(1.0), Inches(8), Inches(4.5),
            )
            _pptx_text(slide, Inches(0.8), Inches(5.7), Inches(8.2), Inches(1.5),
                       section["commentary"][:400], size=10, color=text_dark)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_text(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                       section["title"], size=22, bold=True, color=primary)
            _pptx_text(slide, Inches(0.8), Inches(1.2), Inches(8.2), Inches(5.5),
                       section["commentary"], size=13, color=text_dark)

        if section.get("include_holdings_table") or chart_id in ("contributor_bar", "pie_chart"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _pptx_text(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.5),
                       "Holdings Detail", size=22, bold=True, color=primary)
            _pptx_table(slide, data["holdings"],
                        Inches(0.3), Inches(1.0), Inches(9.4), Inches(5.0), theme)

    prs.save(output_path)
    return output_path
